#!/usr/bin/env python3
"""EIGENNEXUS CxO deck (judge-safe, 9 slides; internal red-team panel omitted). Dark, 16:9.
All numbers cross-checked against results/*.json and results/*.md; chart values are kept
identical to the paper's Fig. 3 (make_coherence_wall_fig.py)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG    = RGBColor(0x0A, 0x0F, 0x1E)
CARD  = RGBColor(0x12, 0x1A, 0x2E)
CARD2 = RGBColor(0x0E, 0x15, 0x27)
RULE  = RGBColor(0x24, 0x32, 0x4C)
INK   = RGBColor(0xFF, 0xFF, 0xFF)
INK2  = RGBColor(0xA9, 0xB4, 0xC8)
MUTED = RGBColor(0x6E, 0x7A, 0x90)
CYAN  = RGBColor(0x4F, 0xD8, 0xEA)
AMBER = RGBColor(0xEE, 0xB2, 0x4C)
GOOD  = RGBColor(0x57, 0xC8, 0x8B)
BAD   = RGBColor(0xE8, 0x6A, 0x6A)
F, FM = "Calibri", "Courier New"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 9


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=0.75, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, font=F,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str): runs = [[(runs, {})]]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(space_after)
        for t, kw in para:
            r = p.add_run(); r.text = t
            r.font.name = kw.get("font", font); r.font.size = Pt(kw.get("size", size))
            r.font.bold = kw.get("bold", bold); r.font.color.rgb = kw.get("color", color)
            if kw.get("italic"): r.font.italic = True
    return tb


def footer(s, n):
    text(s, 0.6, 7.08, 8, 0.3, [[("EIGENNEXUS  ·  CHIMERA-QRC  ·  GIC 2026 Track A", {})]],
         size=9, color=MUTED)
    text(s, 12.1, 7.08, 0.7, 0.3, [[(f"{n:02d} / {TOTAL:02d}", {})]],
         size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def title_row(s, kicker, title):
    text(s, 0.6, 0.42, 12.1, 0.3, [[(kicker.upper(), {})]], size=11, color=CYAN, bold=True)
    text(s, 0.6, 0.72, 12.1, 0.85, [[(title, {})]], size=30, bold=True)


# ---------------------------------------------------------------- 1 · title
s = slide()
box(s, 0, 5.9, 13.333, 1.6, fill=CARD2)
text(s, 0.9, 1.6, 11.5, 0.4, [[("TEAM EIGENNEXUS  ·  GLOBAL INDUSTRY CHALLENGE 2026  ·  QBRAID × MITRE × JONESTRADING", {})]],
     size=12, color=CYAN, bold=True)
text(s, 0.9, 2.05, 11.5, 1.9,
     [[("CHIMERA-QRC", {"size": 54, "bold": True})],
      [("The Honest Quantum Audit", {"size": 30, "color": INK2})]],
     spacing=1.02)
text(s, 0.9, 4.15, 11.5, 0.9,
     [[("Quantum reservoir computing for financial volatility — Track A.  ", {"size": 15, "color": INK2}),
       ("13 QPU campaigns · three vendors · four devices · every claim controlled.", {"size": 15, "bold": True})]])
text(s, 0.9, 6.2, 11.5, 1.0,
     [[("Christian Metzl (Lead / Architect)   ·   Fares Eldibani (Data Science)   ·   Juan Manuel Aguiar Hualde (PhD Physics)", {"size": 12, "color": INK2})],
      [("Every number in this deck traces to a committed artifact and a platform-timestamped job record.", {"size": 11, "color": MUTED, "italic": True})]])

# ---------------------------------------------------------------- 2 · TL;DR
s = slide()
title_row(s, "The instrument first, then the verdict", "What we built — and what it shows")
cards = [
    ("1 command", CYAN, "reproduces the entire audit — offline, no credits",
     "A pre-registered instrument: engine tests + hardware bootstrap + noise fingerprint, re-run in "
     "~2 minutes with the network cut. Every hardware prediction was hash-committed to a prior commit "
     "before the data existed. This machine — not any single number — is the durable asset."),
    ("3 devices", GOOD, "retained quantum signal on real hardware",
     "Trapped-ion (IonQ) and two superconducting chips (IQM Garnet at n=10/12, IQM Emerald) executed our "
     "reservoir with signal intact — each bootstrapped verdict 9.7–23.9σ beyond shot noise, refuting our own "
     "pre-registered predictions under controls. Billing audited to the half-credit across ≈64k credits."),
    ("$0", AMBER, "of measured quantum advantage for vol-forecasting today",
     "Across every fair, pre-registered test: parity at best, never advantage. After Holm nothing is "
     "significant either way and the 95% Model Confidence Set retains all models — a statistical tie, "
     "with the simplest classical model (HAR-X) the rational choice. The avoided-cost result, with receipts."),
]
for i, (stat, col, head, body) in enumerate(cards):
    x = 0.6 + i * 4.15
    box(s, x, 1.85, 3.9, 4.6, fill=CARD, round_=True)
    text(s, x + 0.3, 2.2, 3.3, 1.0, [[(stat, {"size": 40, "bold": True, "color": col})]])
    text(s, x + 0.3, 3.15, 3.3, 0.75, [[(head, {"size": 15, "bold": True})]], spacing=1.05)
    text(s, x + 0.3, 3.95, 3.3, 2.3, [[(body, {"size": 12.5, "color": INK2})]], spacing=1.12)
text(s, 0.6, 6.62, 12.1, 0.35,
     [[("The differentiator is not a claim — it is the audit machinery that makes every claim checkable.", {"size": 13, "italic": True, "color": INK2})]])
footer(s, 2)

# ------------------------------------------------- 3 · the measurement chart
s = slide()
title_row(s, "Four devices, three vendors, one picture", "The coherence wall, measured from both sides")
# bars: (label, raw, limit, signal_bearing, note)
# Identical config list and values to the paper's Fig. 3 (make_coherence_wall_fig.py) so the
# two charts of the same wall cannot disagree. Ticks are instance-matched limits (mean|F_exact|).
bars = [
    ("IonQ Forte-1\nn=8 · ion", 0.104, 0.196, True,  "500 shots"),
    ("IQM Emerald\nn=8 · new gen", 0.169, 0.196, True, ""),
    ("IQM Garnet\nn=10",  0.159, 0.179, True,  ""),
    ("IQM Garnet\nn=12",  0.190, 0.214, True,  ""),
    ("Garnet n=8\nseed-1 (S7)", 0.159, 0.1806, True, "instance limit"),
    ("Garnet n=8\nseed-0", 0.228, 0.196, False, "4 runs 0.222–0.231"),
    ("Rigetti Cep-1\nn=8", 0.223, 0.196, False, "4k rep; 2k was 0.261"),
]
x0, y_base, bw, gap, hmax = 1.05, 5.55, 1.30, 0.34, 3.1  # hmax inches at err=0.28
for i, (lab, raw, lim, sig, note) in enumerate(bars):
    x = x0 + i * (bw + gap)
    h = raw / 0.28 * hmax
    col = GOOD if sig else BAD
    box(s, x, y_base - h, bw, h, fill=col)
    # limit tick
    lh = lim / 0.28 * hmax
    box(s, x - 0.07, y_base - lh - 0.012, bw + 0.14, 0.024, fill=INK)
    text(s, x - 0.1, y_base - h - 0.34, bw + 0.2, 0.3,
         [[(f"{raw:.3f}", {"size": 13, "bold": True, "color": col})]], align=PP_ALIGN.CENTER)
    text(s, x - 0.18, y_base + 0.08, bw + 0.36, 0.62, [[(lab, {"size": 10.5, "color": INK2})]],
         align=PP_ALIGN.CENTER, spacing=0.95)
    if note:
        text(s, x - 0.18, y_base + 0.62, bw + 0.36, 0.3, [[(note, {"size": 9, "color": MUTED})]],
             align=PP_ALIGN.CENTER)
leg_y = 1.78
box(s, 9.05, leg_y + 0.03, 0.24, 0.16, fill=GOOD)
text(s, 9.36, leg_y - 0.04, 1.9, 0.3, [[("signal-bearing", {"size": 11, "color": INK2})]])
box(s, 11.05, leg_y + 0.03, 0.24, 0.16, fill=BAD)
text(s, 11.36, leg_y - 0.04, 1.35, 0.3, [[("scrambled", {"size": 11, "color": INK2})]])
box(s, 9.05, leg_y + 0.33, 0.24, 0.035, fill=INK)
text(s, 9.36, leg_y + 0.21, 3.6, 0.3, [[("size-matched depolarized limit", {"size": 11, "color": INK2})]])
text(s, 0.6, 6.45, 12.1, 0.55,
     [[("Bars: raw feature error vs the exact simulation (4,000 shots; IonQ 500). Below the white tick, the device retains the circuit's signal; "
        "above it, the output is statistically indistinguishable from a fully scrambled state. Every bootstrapped verdict (9 campaigns with committed counts) sits 9.7–23.9σ from its tick.", {"size": 11.5, "color": INK2})]],
     spacing=1.1)
footer(s, 3)

# ------------------------------------------- 4 · controls & self-refutation
s = slide()
title_row(s, "Why these results survive hostile review", "We bought controls — and refuted our own predictions")
rows = [
    ("Same-session anchor", "n=8 and n=12 interleaved on one chip, matching calibration fingerprints",
     "The effect is real for this instance — day-drift excluded by a rule committed before the data existed; S7 later localized it to the seed-0 instance, not size.", GOOD),
    ("Same-window two-chip pair", "Garnet and Emerald on the clock simultaneously, first jobs the same second",
     "The generation effect is temporally controlled: newer chip signal-bearing at the size the older one scrambles.", GOOD),
    ("4k-shot replication", "Rigetti re-run at double shots on a second day",
     "Day-scale drift measured (~0.04) — larger than shot noise; regime claims, not point values, are the currency.", AMBER),
]
y = 1.9
for head, what, why, col in rows:
    box(s, 0.6, y, 12.1, 1.35, fill=CARD, round_=True)
    text(s, 0.95, y + 0.18, 3.1, 1.0, [[(head, {"size": 15, "bold": True, "color": col})]], spacing=1.02)
    text(s, 4.25, y + 0.18, 4.1, 1.05, [[(what, {"size": 12, "color": INK2})]], spacing=1.08)
    text(s, 8.55, y + 0.18, 3.85, 1.05, [[(why, {"size": 12})]], spacing=1.08)
    y += 1.53
text(s, 0.6, y + 0.03, 12.1, 0.5,
     [[("Scorecard of our own pre-registered statements:  ", {"size": 12.5, "color": INK2}),
       ("S1 · S2 · S3b refuted · S7 cross-seed refuted · S4 held", {"size": 12.5, "bold": True}),
       (" — four of our own pre-registered predictions falsified by controlled measurement, committed first, reported as measured.", {"size": 12.5, "color": INK2})]],
     spacing=1.05)
footer(s, 4)

# ---------------------------------------------------- 5 · business impact
s = slide()
title_row(s, "What this means commercially", "Five implications a board can act on")
impl = [
    ("01", "No quantum spend on volatility forecasting today",
     "Parity, never advantage, across crisis and calm — the avoided-cost result, in writing, with receipts."),
    ("02", "Efficiency is a negative too",
     "On like-for-like accuracy (lower-is-better RMSE), the quantum reservoir saturates worst. Parity is the only surviving edge."),
    ("03", "The readiness buffer just shrank",
     "Signal-bearing execution arrived on three devices — the pre-registered upside (S3 failing high) that we said would materially upgrade the outlook, and it did. Move from ignore to monitor."),
    ("04", "Spec-sheet procurement is unreliable",
     "Gate-count heuristics mis-ranked both instances we measured on metal (n=2). Qualify hardware per workload, empirically — the mechanism is still open (H-EMBED pre-registered, unrun)."),
    ("05", "Certainty is the product",
     "A pre-registered audit converts hype into policy at a known, small cost — and it is repeatable on demand."),
]
y = 1.85
for num, head, body in impl:
    box(s, 0.6, y, 12.1, 0.92, fill=CARD2, round_=True)
    text(s, 0.92, y + 0.16, 0.75, 0.6, [[(num, {"size": 22, "bold": True, "color": CYAN, "font": FM})]])
    text(s, 1.85, y + 0.13, 4.6, 0.66, [[(head, {"size": 14.5, "bold": True})]], spacing=1.0,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.65, y + 0.13, 5.75, 0.72, [[(body, {"size": 11.5, "color": INK2})]], spacing=1.05,
         anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
footer(s, 5)

# ---------------------------------------------------- 6 · trust architecture
s = slide()
title_row(s, "The audit machinery", "Trust is engineered, not asserted")
steps = [
    ("Pre-register", "Predictions, budgets, abort and decision rules committed to the repo before execution — amendments timestamped per campaign."),
    ("Tag every job", "Each QPU job in every funded campaign embeds the repo commit hash in qBraid's records: a hash-preimage proof that predictions predate data."),
    ("Audit every credit", "60,698.25 org credits across 10 funded campaigns (64,048.25 incl. a disclosed personal-account anomaly), each at or under its pre-approved reservation; organizer-grade ledger with per-job IDs."),
    ("Re-derive everything", "Bootstrap from committed raw counts re-derives the nine campaigns with committed counts (9.7–23.9σ)."),
]
for i, (head, body) in enumerate(steps):
    x = 0.6 + i * 3.22
    box(s, x, 1.95, 2.95, 3.3, fill=CARD, round_=True)
    text(s, x + 0.25, 2.2, 2.45, 0.5, [[(f"{i+1}", {"size": 26, "bold": True, "color": CYAN, "font": FM})]])
    text(s, x + 0.25, 2.85, 2.45, 0.6, [[(head, {"size": 15, "bold": True})]])
    text(s, x + 0.25, 3.5, 2.45, 1.6, [[(body, {"size": 11.5, "color": INK2})]], spacing=1.12)
box(s, 0.6, 5.55, 12.1, 1.15, fill=CARD2, round_=True)
text(s, 0.95, 5.75, 11.4, 0.8,
     [[("Side effect of auditing the platform this hard: ", {"size": 12.5, "color": INK2}),
       ("two reproducible platform findings documented with reproduction bundles (plus one billing reconstruction, explicitly unconfirmed)", {"size": 12.5, "bold": True}),
       (" (silent negative-angle gate loss; a simulator anomaly; a billing-context fallback) — finding 3 not yet vendor-confirmed and labeled as our reconstruction. "
        "The audit made the ecosystem better while competing on it.", {"size": 12.5, "color": INK2})]],
     spacing=1.12)
footer(s, 6)


# ---------------------------------------------------- 7 · roadmap
s = slide()
title_row(s, "Committed next chapter", "Pre-registered research · commercial ladder")
box(s, 0.6, 1.9, 5.9, 4.6, fill=CARD, round_=True)
text(s, 0.95, 2.15, 5.2, 0.4, [[("PRE-REGISTERED, FALSIFIABLE, COSTED", {"size": 11, "bold": True, "color": CYAN})]])
road = [
    ("S5", "Hardware-in-the-loop forecast on the signal-bearing Emerald — degradation band computed before any run; 'does not beat HAR-X' committed."),
    ("S6", "Hardware-native dissipative reservoirs — the simulation's +60% autonomous-VPT mechanism, tested where damping is free."),
    ("S7 ✓", "EXECUTED (cross-seed control): an independent seed-1 n=8 instance is signal-bearing while seed-0 scrambles — refutes a size law (in the paper). The embedding arm (H-EMBED) remains pre-registered for future work."),
]
y = 2.6
for tag, body in road:
    text(s, 0.95, y, 0.7, 0.4, [[(tag, {"size": 15, "bold": True, "color": CYAN, "font": FM})]])
    text(s, 1.75, y, 4.5, 1.05, [[(body, {"size": 12, "color": INK2})]], spacing=1.1)
    y += 1.24
box(s, 6.85, 1.9, 5.85, 4.6, fill=CARD, round_=True)
text(s, 7.2, 2.15, 5.2, 0.4, [[("COMMERCIAL LADDER", {"size": 11, "bold": True, "color": AMBER})]])
ladder = [
    ("Audit-as-a-service", "Vendor-neutral hardware qualification for a named workload — the rigged-to-be-honest exam, priced per campaign."),
    ("Tripwire monitoring", "Annual re-run of the sentinel harness: a few thousand credits to know when the answer changes."),
    ("Qualification consulting", "Spec sheets demonstrably mislead; per-workload empirical qualification is now the defensible standard."),
]
y = 2.6
for head, body in ladder:
    text(s, 7.2, y, 5.2, 0.35, [[(head, {"size": 13.5, "bold": True})]])
    text(s, 7.2, y + 0.38, 5.2, 0.75, [[(body, {"size": 11.5, "color": INK2})]], spacing=1.08)
    y += 1.24
footer(s, 7)

# ---------------------------------------------------- 8 · audit-as-a-service
s = slide()
title_row(s, "The commercial artifact", "We sell trust in a quantum yes/no — evidenced by an instrument that falsifies itself")
# offering (left)
box(s, 0.6, 1.9, 5.0, 4.6, fill=CARD, round_=True)
text(s, 0.92, 2.15, 4.4, 0.4, [[("THE OFFERING", {"size": 11, "bold": True, "color": CYAN})]])
text(s, 0.92, 2.6, 4.4, 3.6,
     [[("What the client buys is a trustworthy yes/no: does a given QPU retain signal on their OWN circuit, before they commit budget? — credible because the same instrument publishes negatives about itself.", {"size": 13, "bold": True})],
      [("• pre-registered predictions + decision rules, committed before any run", {"size": 12, "color": INK2})],
      [("• controlled measurement across vendors (same-session / same-window)", {"size": 12, "color": INK2})],
      [("• a report of what each QPU can and cannot do — every number re-derivable from raw data", {"size": 12, "color": INK2})],
      [("• proof it works: our own program qualified 5 configs signal-bearing, 2 scrambled — on the SAME circuit (chart, slide 3); a ~$70 single-device test run before you commit to a build", {"size": 12, "color": INK2})]],
     spacing=1.15, space_after=8)
# pricing (middle) — REAL cost basis from our measured spend
box(s, 5.8, 1.9, 3.4, 4.6, fill=CARD2, round_=True)
text(s, 6.05, 2.15, 2.9, 0.4, [[("COST BASIS (MEASURED)", {"size": 11, "bold": True, "color": AMBER})]])
text(s, 6.05, 2.62, 2.95, 3.7,
     [[("Single-device qualification", {"size": 12.5, "bold": True})],
      [("~6,800 credits (~$70 QPU) → one report", {"size": 11, "color": INK2})],
      [("Multi-vendor program", {"size": 12.5, "bold": True})],
      [("~64k credits (~$640 QPU) → 4 devices, controlled", {"size": 11, "color": INK2})],
      [("These are our actual measured QPU costs (CREDIT_BUDGET.md).", {"size": 10.5, "color": MUTED, "italic": True})],
      [("Service fee is set per engagement; the value is the method + analysis, not the compute.", {"size": 10.5, "color": MUTED, "italic": True})]],
     spacing=1.12, space_after=7)
# target segments (right) — profiles, NOT named customers
box(s, 9.4, 1.9, 3.3, 4.6, fill=CARD, round_=True)
text(s, 9.65, 2.15, 2.85, 0.4, [[("WHO BUYS FIRST", {"size": 11, "bold": True, "color": CYAN})]])
text(s, 9.65, 2.62, 2.85, 3.7,
     [[("BEACHHEAD — a buy-side risk / quant desk", {"size": 11.5, "bold": True, "color": INK2})],
      [("about to fund a quantum pilot: we tell them which device (if any) survives THEIR circuit — before they build.", {"size": 11, "color": INK2})],
      [("EXPANSION", {"size": 11.5, "bold": True, "color": INK2})],
      [("Hardware vendors (independent benchmarks) · CTO quantum-readiness diligence · public-sector assurance.", {"size": 11, "color": INK2})]],
     spacing=1.12, space_after=10)
text(s, 0.6, 6.65, 12.1, 0.35,
     [[("Honesty note: segments are target profiles, not signed customers. ", {"size": 11, "italic": True, "color": AMBER}),
       ("Pilot LOIs are to be secured — none are claimed here.", {"size": 11, "italic": True, "color": INK2})]])
footer(s, 8)

# ---------------------------------------------------- 9 · close
s = slide()
box(s, 0, 2.3, 13.333, 2.6, fill=CARD2)
text(s, 0.9, 2.75, 11.5, 1.6,
     [[("Truth, timed and priced.", {"size": 40, "bold": True})],
      [("An honest negative, a controlled surprise, and the machinery to keep both current.", {"size": 16, "color": INK2})]],
     spacing=1.15)
text(s, 0.9, 5.35, 11.5, 1.3,
     [[("Team EIGENNEXUS — Christian Metzl · Fares Eldibani · Juan Manuel Aguiar Hualde", {"size": 13, "bold": True})],
      [("Full evidence: submission repository (paper, findings, ledgers, job IDs, reproduction in one command).", {"size": 11.5, "color": INK2})],
      [("Every number here is re-derivable from committed data; the pre-registered claims are additionally hash-committed in advance.", {"size": 11.5, "color": MUTED, "italic": True})]],
     spacing=1.2, space_after=5)
footer(s, 9)

OUT = "docs/EIGENNEXUS_CxO_Deck.pptx"
prs.save(OUT)
print("wrote", OUT)
